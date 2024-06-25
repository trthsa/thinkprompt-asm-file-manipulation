import os
import json
from deep_translator import GoogleTranslator
from docx import Document
from docx.shared import Pt, RGBColor
import fitz  # PyMuPDF
from pptx import Presentation

class FileManipulator:
    def __init__(self):
        pass

    def _ensure_folder_exists(self, folder_path):
        """
        Ensure the given folder exists, create it if it does not.

        :param folder_path: Path to the folder.
        """
        if not os.path.exists(folder_path):
            os.makedirs(folder_path)

    def _save_metadata(self, metadata, output_folder):
        """
        Save metadata to a JSON file in the specified folder.

        :param metadata: Metadata to be saved.
        :param output_folder: Path to the folder where metadata will be saved.
        """
        with open(os.path.join(output_folder, "metadata.json"), "w") as meta_file:
            json.dump(metadata, meta_file, indent=4)

    def _load_metadata(self, output_folder):
        """
        Load metadata from a JSON file in the specified folder.

        :param output_folder: Path to the folder containing the metadata.json file.
        :return: Loaded metadata.
        """
        with open(os.path.join(output_folder, "metadata.json"), "r") as meta_file:
            return json.load(meta_file)

    def _process_text(self, text, processor=None):
        """
        Process the given text using the specified processor function.

        :param text: Text to be processed.
        :param processor: Optional function to process the text.
        :return: Processed text.
        """
        return processor(text) if processor else text

    def extract_text_images_from_pdf(self, pdf_path, output_folder):
        """
        Extract text and images from a PDF and save metadata to a specified folder.

        :param pdf_path: Path to the input PDF file.
        :param output_folder: Path to the folder where output metadata and images will be saved.
        """
        self._ensure_folder_exists(output_folder)

        doc = fitz.open(pdf_path)
        metadata = {"pages": []}

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)

            # Extract text with styling information
            page_blocks = page.get_text("dict")["blocks"]
            text_metadata = []

            for block in page_blocks:
                if block["type"] == 0:  # Text block
                    for line in block["lines"]:
                        for span in line["spans"]:
                            bbox = span["bbox"]
                            text_meta = {
                                "text": span["text"],
                                "font_size": span["size"],
                                "font_color": span["color"],
                                "font_name": span["font"],
                                "bbox": bbox,
                                "bold": span["flags"] & 1 != 0 or ("bold" in span["font"].lower()),  
                                "italic": span["flags"] & 2 != 0 or ("italic" in span["font"].lower()), 
                            }
                            text_metadata.append(text_meta)

            # Extract images
            images = page.get_images(full=True)
            image_metadata = []
            for img_index, img in enumerate(images):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                image_path = os.path.join(output_folder, f"page_{page_num + 1}_image_{img_index + 1}.{image_ext}")
                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)
                image_metadata.append({
                    "index": img_index,
                    "path": image_path,
                    "ext": image_ext
                })

            metadata["pages"].append({
                "page_num": page_num + 1,
                "text_instances": text_metadata,
                "images": image_metadata
            })

        self._save_metadata(metadata, output_folder)

    def extract_text_images_from_docx(self, docx_path, output_folder):
        """
        Extract text and images from a DOCX file and save metadata to a specified folder.

        :param docx_path: Path to the input DOCX file.
        :param output_folder: Path to the folder where output metadata and images will be saved.
        """
        self._ensure_folder_exists(output_folder)

        doc = Document(docx_path)
        metadata = {"paragraphs": [], "images": [], "tables": []}

        # Paragraph extraction
        for para_num, para in enumerate(doc.paragraphs):
            para_data = {
                "para_num": para_num + 1,
                "text": para.text,
                "runs": []
            }
            for run in para.runs:
                run_data = {
                    "text": run.text,
                    "bold": run.bold,
                    "italic": run.italic,
                    "font_name": run.font.name,
                    "font_size": run.font.size.pt if run.font.size else None,
                    "color": run.font.color.rgb if run.font.color else None
                }
                para_data["runs"].append(run_data)
            metadata["paragraphs"].append(para_data)

        # Image extraction
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img = rel.target_part.blob
                img_ext = rel.target_part.content_type.split('/')[-1]
                img_path = os.path.join(output_folder, f"image_{rel.target_ref.split('/')[-1]}.{img_ext}")
                with open(img_path, "wb") as img_file:
                    img_file.write(img)
                metadata["images"].append({
                    "path": img_path,
                    "ext": img_ext
                })

        # Table extraction with text styles
        for table_num, table in enumerate(doc.tables):
            table_data = {"table_num": table_num + 1, "rows": []}
            for row_num, row in enumerate(table.rows):
                row_data = {"row_num": row_num + 1, "cells": []}
                for cell in row.cells:
                    cell_data = {"paragraphs": []}
                    for para in cell.paragraphs:
                        para_data = {
                            "text": para.text,
                            "runs": []
                        }
                        for run in para.runs:
                            run_data = {
                                "text": run.text,
                                "bold": run.bold,
                                "italic": run.italic,
                                "font_name": run.font.name,
                                "font_size": run.font.size.pt if run.font.size else None,
                                "color": run.font.color.rgb if run.font.color else None
                            }
                            para_data["runs"].append(run_data)
                        cell_data["paragraphs"].append(para_data)
                    row_data["cells"].append(cell_data)
                table_data["rows"].append(row_data)
            metadata["tables"].append(table_data)

        self._save_metadata(metadata, output_folder)

    def convert_text_to_uppercase(self, file_path, output_path, file_type):
        """
        Convert all text in a file (PDF or DOCX) to uppercase and save the modified file.

        :param file_path: Path to the input file.
        :param output_path: Path to save the modified file.
        :param file_type: Type of the input file ('pdf' or 'docx').
        """
        if file_type == "pdf":
            self.recreate_pdf(file_path, output_path, lambda text: text.upper())
        elif file_type == "docx":
            self.recreate_docx(file_path, output_path, lambda text: text.upper())

    def recreate_docx(self, output_folder, output_path, content_processor=None):
        """
        Recreate a DOCX file from extracted metadata and apply a content processor function to the text.

        :param output_folder: Path to the folder containing the metadata.json file.
        :param output_path: Path to save the recreated DOCX file.
        :param content_processor: Optional function to process the text content.
        """
        metadata = self._load_metadata(output_folder)

        doc = Document()

        # Paragraphs recreation
        for para_meta in metadata["paragraphs"]:
            para = doc.add_paragraph()
            for run_meta in para_meta["runs"]:
                text = self._process_text(run_meta["text"], content_processor)
                run = para.add_run(text)
                run.bold = run_meta.get("bold", False)
                run.italic = run_meta.get("italic", False)
                if run_meta.get("font_size"):
                    run.font.size = Pt(run_meta["font_size"])
                if run_meta.get("color"):
                    run.font.color.rgb = RGBColor(*run_meta["color"])
                run.font.name = run_meta.get("font_name")

        # Images recreation
        for img_meta in metadata["images"]:
            with open(img_meta["path"], "rb") as img_file:
                doc.add_picture(img_file, width=Pt(300))

        # Tables recreation
        for table_meta in metadata["tables"]:
            table = doc.add_table(rows=0, cols=len(table_meta["rows"][0]["cells"]) if table_meta["rows"] else 0)
            for row_meta in table_meta["rows"]:
                row = table.add_row()
                for cell_meta, cell in zip(row_meta["cells"], row.cells):
                    cell._element.clear_content()  # Clear existing cell content
                    for para_meta in cell_meta["paragraphs"]:
                        para = cell.add_paragraph()
                        for run_meta in para_meta["runs"]:
                            text = self._process_text(run_meta["text"], content_processor)
                            run = para.add_run(text)
                            run.bold = run_meta.get("bold", False)
                            run.italic = run_meta.get("italic", False)
                            if run_meta.get("font_size"):
                                run.font.size = Pt(run_meta["font_size"])
                            if run_meta.get("color"):
                                run.font.color.rgb = RGBColor(*run_meta["color"])
                            run.font.name = run_meta.get("font_name")

        doc.save(output_path)

    def recreate_pdf(self, output_folder, output_path, content_processor=None):
        """
        Recreate a PDF file from extracted metadata and apply a content processor function to the text.

        :param output_folder: Path to the folder containing the metadata.json file.
        :param output_path: Path to save the recreated PDF file.
        :param content_processor: Optional function to process the text content.
        """
        metadata = self._load_metadata(output_folder)

        doc = fitz.open()

        for page_meta in metadata["pages"]:
            page = doc.new_page()

            for text_meta in page_meta["text_instances"]:
                bbox = fitz.Rect(text_meta["bbox"])
                r = (text_meta["font_color"] >> 16) & 0xff
                g = (text_meta["font_color"] >> 8) & 0xff
                b = text_meta["font_color"] & 0xff
                color = [r/255, g/255, b/255]
                page.insert_text((bbox.x0, bbox.y0),
                                 self._process_text(text_meta["text"], content_processor),
                                 fontsize=text_meta["font_size"], color=color)
            for img_meta in page_meta["images"]:
                image_rect = fitz.Rect(72, 72 + img_meta["index"] * 200, 300, 400 + img_meta["index"] * 200)
                page.insert_image(image_rect, filename=img_meta["path"])

        doc.save(output_path)

    def translate_text_in_pptx(self, input_pptx_path, output_pptx_path,from_lang='en',to_lang='vi'):
        """
        Translate all text in a PPTX file to Vietnamese and append the translated text under the original text.

        :param input_pptx_path: Path to the input PPTX file.
        :param output_pptx_path: Path to save the modified PPTX file.
        """
        prs = Presentation(input_pptx_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        translated_text = GoogleTranslator(source=from_lang, target=to_lang).translate(run.text)
                        run.text = run.text + "\n" + translated_text

        prs.save(output_pptx_path)


if __name__ == "__main__":
    manipulator = FileManipulator()

    # # Extract text and images from PDF
    manipulator.extract_text_images_from_pdf("mock_samples/pdf_mock_file.pdf", "results/pdf_mock_file_output")

    # # Extract text and images from DOCX
    manipulator.extract_text_images_from_docx("mock_samples/docx_mock_file.docx", "results/docx_mock_file_output")

    # # Convert text to uppercase and compile into new DOCX
    manipulator.convert_text_to_uppercase("results/docx_mock_file_output", "results/uppercase_docx.docx", file_type="docx")

    # # Convert text to uppercase and compile into new PDF
    manipulator.convert_text_to_uppercase("results/pdf_mock_file_output", "results/uppercase_pdf.pdf", file_type="pdf")

    # # Recreate DOCX from extracted data
    manipulator.recreate_docx("results/docx_mock_file_output", "recreated_docx.docx")

    # # Recreate PDF from extracted data
    manipulator.recreate_pdf("results/pdf_mock_file_output", "recreated_pdf.pdf")

    # Translate all text in the file to Vietnamese and append the translated text under the original text in slides.
    # manipulator.translate_text_in_pptx("mock_samples/Networking.pptx", "results/translated_pptx.pptx",from_lang='en',to_lang='vi')
