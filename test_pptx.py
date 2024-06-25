from deep_translator import GoogleTranslator
from pptx import Presentation

def change_text_in_pptx(input_pptx_path, output_pptx_path):
    """
    Change text in a PPTX file.

    :param input_pptx_path: Path to the input PPTX file.
    :param output_pptx_path: Path to save the modified PPTX file.
    :param text_replacements: Dictionary with keys as old text and values as new text.
    """
    prs = Presentation(input_pptx_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = run.text + "\n" + GoogleTranslator(source='en', target='vi').translate(run.text) 

    prs.save(output_pptx_path)

# Example usage
input_pptx_path = "Networking.pptx"
output_pptx_path = "modified_example.pptx"
change_text_in_pptx(input_pptx_path, output_pptx_path)
